#!/usr/bin/env python3
"""
Generate hubbleAI Sales Team Presentation with Aperam branding
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Aperam brand colors
PURPLE = RGBColor(91, 45, 142)       # #5B2D8E
ORANGE = RGBColor(232, 93, 4)        # #E85D04
GREEN = RGBColor(46, 125, 50)        # #2E7D32
DARK_GRAY = RGBColor(64, 64, 64)     # #404040
LIGHT_GRAY = RGBColor(245, 245, 245) # #F5F5F5
WHITE = RGBColor(255, 255, 255)

def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide with purple header bar"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar at top
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(2.5)
    )
    set_shape_fill(header, PURPLE)
    header.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.8),
        Inches(11), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.9),
            Inches(11), Inches(0.6)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.LEFT

    return slide

def add_content_slide(prs, title, content_items, accent_color=PURPLE):
    """Add a content slide with colored left border boxes"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Thin purple line at top
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(12.333), Pt(3)
    )
    set_shape_fill(line, PURPLE)
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.7),
        Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Content boxes with colored left border
    y_pos = 1.7
    for i, item in enumerate(content_items):
        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Pt(6), Inches(0.9)
        )
        bar_color = accent_color if isinstance(accent_color, RGBColor) else [PURPLE, ORANGE, GREEN][i % 3]
        set_shape_fill(bar, bar_color)
        bar.line.fill.background()

        # Content box background
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(y_pos),
            Inches(12), Inches(0.9)
        )
        set_shape_fill(box, LIGHT_GRAY)
        box.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(y_pos + 0.15),
            Inches(11.5), Inches(0.7)
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        if isinstance(item, tuple):
            # Bold header + description
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item[0]
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = DARK_GRAY

            p2 = tf.add_paragraph()
            p2.text = item[1]
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK_GRAY
        else:
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY

        y_pos += 1.1

    return slide

def add_comparison_slide(prs, title, before_items, after_items):
    """Add a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Thin purple line at top
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(12.333), Pt(3)
    )
    set_shape_fill(line, PURPLE)
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.7),
        Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Before column header
    before_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.7),
        Inches(5.8), Inches(0.6)
    )
    set_shape_fill(before_header, ORANGE)
    before_header.line.fill.background()

    bh_text = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.75),
        Inches(5.8), Inches(0.5)
    )
    tf = bh_text.text_frame
    p = tf.paragraphs[0]
    p.text = "BEFORE (Traditional)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # After column header
    after_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.8), Inches(1.7),
        Inches(5.8), Inches(0.6)
    )
    set_shape_fill(after_header, GREEN)
    after_header.line.fill.background()

    ah_text = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.75),
        Inches(5.8), Inches(0.5)
    )
    tf = ah_text.text_frame
    p = tf.paragraphs[0]
    p.text = "AFTER (hubbleAI)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Before items
    y_pos = 2.5
    for item in before_items:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(5.8), Inches(0.7)
        )
        set_shape_fill(box, LIGHT_GRAY)
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.15),
            Inches(5.4), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY

        y_pos += 0.85

    # After items
    y_pos = 2.5
    for item in after_items:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.8), Inches(y_pos),
            Inches(5.8), Inches(0.7)
        )
        set_shape_fill(box, LIGHT_GRAY)
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(y_pos + 0.15),
            Inches(5.4), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "✓  " + item
        p.font.size = Pt(18)
        p.font.color.rgb = GREEN

        y_pos += 0.85

    return slide

def add_metrics_slide(prs, title, metrics):
    """Add a slide with metric boxes"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Thin purple line at top
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(12.333), Pt(3)
    )
    set_shape_fill(line, PURPLE)
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.7),
        Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Metric boxes (2 rows of 2)
    colors = [PURPLE, ORANGE, GREEN, RGBColor(220, 180, 50)]
    positions = [
        (0.5, 1.8), (6.6, 1.8),
        (0.5, 4.0), (6.6, 4.0)
    ]

    for i, (metric, value) in enumerate(metrics[:4]):
        x, y = positions[i]
        color = colors[i]

        # Box outline
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(5.8), Inches(1.8)
        )
        box.fill.background()
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # Value
        val_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.3),
            Inches(5.8), Inches(0.9)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 1.2),
            Inches(5.8), Inches(0.5)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_diagram_slide(prs, title):
    """Add the 'How it works' diagram slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Thin purple line at top
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(12.333), Pt(3)
    )
    set_shape_fill(line, PURPLE)
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.7),
        Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Three boxes: Learn -> Predict -> Compare
    steps = [
        ("1. LEARN", "AI studies 1+ year\nof historical data", PURPLE),
        ("2. PREDICT", "Generates forecasts\nfor next 8 weeks", ORANGE),
        ("3. COMPARE", "Shows performance\nvs traditional method", GREEN),
    ]

    x_positions = [0.8, 4.8, 8.8]

    for i, (header, desc, color) in enumerate(steps):
        x = x_positions[i]

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.8),
            Inches(3.5), Inches(2.2)
        )
        set_shape_fill(box, color)
        box.line.fill.background()

        # Header
        h_box = slide.shapes.add_textbox(
            Inches(x), Inches(2.0),
            Inches(3.5), Inches(0.6)
        )
        tf = h_box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        d_box = slide.shapes.add_textbox(
            Inches(x), Inches(2.6),
            Inches(3.5), Inches(1.2)
        )
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Arrow between boxes
        if i < 2:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + 3.6), Inches(2.7),
                Inches(0.9), Inches(0.5)
            )
            set_shape_fill(arrow, DARK_GRAY)
            arrow.line.fill.background()

    # Bottom text - what makes it smart
    smart_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.3),
        Inches(12), Inches(0.5)
    )
    tf = smart_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What makes it smart:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    points = [
        "✓  Learns seasonal patterns (month-end spikes, holiday dips)",
        "✓  Provides a range: best case / likely / worst case",
        "✓  Runs automatically every week — no manual effort"
    ]

    y = 4.9
    for point in points:
        p_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(y),
            Inches(12), Inches(0.4)
        )
        tf = p_box.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        y += 0.45

    return slide

def add_takeaway_slide(prs, title, main_message, points):
    """Add final takeaway slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(2.0)
    )
    set_shape_fill(header, PURPLE)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),
        Inches(12), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Main message
    msg_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12), Inches(1.0)
    )
    tf = msg_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_message
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # Points with checkmarks
    y = 3.7
    for point in points:
        # Green checkmark box
        check_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(y),
            Inches(10), Inches(0.6)
        )
        check_box.fill.background()
        check_box.line.color.rgb = LIGHT_GRAY
        check_box.line.width = Pt(1)

        check = slide.shapes.add_textbox(
            Inches(1.7), Inches(y + 0.1),
            Inches(0.5), Inches(0.5)
        )
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN

        text = slide.shapes.add_textbox(
            Inches(2.3), Inches(y + 0.12),
            Inches(9), Inches(0.5)
        )
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY

        y += 0.75

    return slide

def add_thank_you_slide(prs):
    """Add closing slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    set_shape_fill(bg, PURPLE)
    bg.line.fill.background()

    # Thank you text
    title_box = slide.shapes.add_textbox(
        Inches(0), Inches(2.5),
        Inches(13.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Questions
    q_box = slide.shapes.add_textbox(
        Inches(0), Inches(4.0),
        Inches(13.333), Inches(0.8)
    )
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER

    # hubbleAI team
    team_box = slide.shapes.add_textbox(
        Inches(0), Inches(5.5),
        Inches(13.333), Inches(0.6)
    )
    tf = team_box.text_frame
    p = tf.paragraphs[0]
    p.text = "hubbleAI Team"
    p.font.size = Pt(20)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "hubbleAI",
        "Smarter Cash Flow Forecasting with AI"
    )

    # Slide 2: The Challenge
    add_content_slide(
        prs,
        "The Challenge We Solved",
        [
            ("Limited visibility", "Traditional forecasts only looked 4 weeks ahead"),
            ("Manual guesswork", "Hard to predict complex cash flows accurately"),
            ("Reactive planning", "Short horizons meant last-minute scrambling"),
            ("Over-budgeting", "Uncertainty led to holding excess cash 'just in case'"),
        ]
    )

    # Slide 3: The Solution (comparison)
    add_comparison_slide(
        prs,
        "Our Solution",
        [
            "4 weeks of visibility",
            "Single-point forecasts",
            "Manual weekly updates",
            "No performance tracking",
        ],
        [
            "8 weeks of visibility",
            "Range: best / likely / worst case",
            "Automated weekly predictions",
            "Transparent accuracy metrics",
        ]
    )

    # Slide 4: How it works
    add_diagram_slide(prs, "How It Works")

    # Slide 5: Results
    add_metrics_slide(
        prs,
        "Results We've Seen",
        [
            ("FORECAST VISIBILITY", "2x"),
            ("WIN RATE VS TRADITIONAL", "75-90%"),
            ("FORECAST ACCURACY", "5-22%"),
            ("RISK COVERAGE", "Best / Likely / Worst"),
        ]
    )

    # Slide 6: Key Takeaway
    add_takeaway_slide(
        prs,
        "Key Takeaway",
        "AI doesn't replace expertise — it extends it.",
        [
            "Extends forecast visibility from 4 to 8 weeks",
            "Learns patterns humans might miss",
            "Always compared against baseline to prove value",
            "Could similar approaches help your team?",
        ]
    )

    # Slide 7: Thank You
    add_thank_you_slide(prs)

    # Save
    output_path = "/home/user/hubbleAI/docs/hubbleAI_Sales_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
